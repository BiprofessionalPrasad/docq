"""Document text extraction for common formats. Fast and dependency-light where possible."""

from __future__ import annotations
import csv
import json
import re
from html.parser import HTMLParser
from pathlib import Path
from typing import Optional

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook


class _TextExtractor(HTMLParser):
    """Very lightweight HTML -> text (no external bs4)."""
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self.skip = False

    def handle_starttag(self, tag: str, attrs: list[tuple[str, Optional[str]]]) -> None:
        tag = tag.lower()
        if tag in {"script", "style", "noscript"}:
            self.skip = True
        elif tag in {"br", "p", "div", "li", "h1", "h2", "h3", "h4", "h5", "h6", "tr"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in {"script", "style", "noscript"}:
            self.skip = False
        elif tag in {"p", "div", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self.skip:
            self.parts.append(data)

    def get_text(self) -> str:
        text = "".join(self.parts)
        # collapse excessive whitespace
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()


def extract_text(path: Path) -> str:
    """Extract plain text from a supported document. Returns '' for unsupported or unreadable."""
    ext = path.suffix.lower()
    try:
        if ext in {".txt", ".md", ".markdown", ".rst", ".log", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".go", ".rs", ".rb", ".sh", ".bat", ".ps1", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".sql", ".r"}:
            return path.read_text(encoding="utf-8", errors="replace")

        if ext in {".json", ".jsonl"}:
            data = path.read_text(encoding="utf-8", errors="replace")
            try:
                obj = json.loads(data)
                return json.dumps(obj, indent=2, ensure_ascii=False)
            except Exception:
                return data

        if ext in {".html", ".htm"}:
            parser = _TextExtractor()
            parser.feed(path.read_text(encoding="utf-8", errors="replace"))
            return parser.get_text()

        if ext == ".csv":
            rows = []
            with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
                reader = csv.reader(f)
                for i, row in enumerate(reader):
                    if i > 200:  # safety for huge CSVs
                        rows.append("... (truncated)")
                        break
                    rows.append(" | ".join(str(c) for c in row))
            return "\n".join(rows)

        if ext == ".pdf":
            reader = PdfReader(str(path))
            pages = []
            for i, page in enumerate(reader.pages):
                t = page.extract_text() or ""
                if t.strip():
                    pages.append(f"[Page {i+1}]\n{t.strip()}")
            return "\n\n".join(pages)

        if ext == ".docx":
            doc = DocxDocument(str(path))
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text.strip())
            # tables too
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
            return "\n".join(parts)

        if ext in {".pptx", ".ppt"}:
            prs = Presentation(str(path))
            slides = []
            for si, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        texts.append(shape.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            if any(cells):
                                texts.append(" | ".join(cells))
                if texts:
                    slides.append(f"[Slide {si}]\n" + "\n".join(texts))
            return "\n\n".join(slides)

        if ext in {".xlsx", ".xls"}:
            wb = load_workbook(str(path), read_only=True, data_only=True)
            sheets = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for ri, row in enumerate(ws.iter_rows(values_only=True)):
                    if ri > 300:
                        rows.append("... (truncated)")
                        break
                    vals = [str(v) if v is not None else "" for v in row]
                    if any(vals):
                        rows.append(" | ".join(vals))
                if rows:
                    sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
            wb.close()
            return "\n\n".join(sheets)

    except Exception as e:
        # Never crash indexing on bad file; log would be nice but for now silent + empty
        return f"[Extraction error for {path.name}: {e}]"

    return ""
